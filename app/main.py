"""
HydroQuote AI - Main Application Entry Point
FastAPI application with secure configuration management
"""
from fastapi import FastAPI, Depends, HTTPException, status, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
import logging
from typing import Dict, Any, Union
import PyPDF2
import docx
from pptx import Presentation
import io
import re
import random
from datetime import datetime
from docx.shared import Pt

from app.core.config import get_settings, Settings

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


def create_application() -> FastAPI:
    """
    Application factory pattern
    Creates and configures the FastAPI application
    """
    settings = get_settings()
    
    # Create FastAPI app with conditional docs
    app = FastAPI(
        title=settings.app_name,
        version=settings.app_version,
        description="Database-free LLM-powered hydro turbine quotation system",
        docs_url="/docs" if settings.enable_swagger_docs else None,
        redoc_url="/redoc" if settings.enable_swagger_docs else None,
    )
    
    # Configure CORS
    app.add_middleware(
        CORSMiddleware,
        allow_origins=settings.cors_origins,
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )
    
    # Log startup information (without sensitive data)
    logger.info(f"Starting {settings.app_name} v{settings.app_version}")
    logger.info(f"Environment: {settings.app_env}")
    logger.info(f"Log Level: {settings.log_level}")
    logger.info(f"API Port: {settings.api_port}")
    logger.info(f"Watsonx Model: {settings.watsonx_model}")
    logger.info(f"CORS Origins: {settings.cors_origins}")
    
    # NEVER log sensitive information
    # ❌ DON'T: logger.info(f"API Key: {settings.watsonx_api_key}")
    # ✅ DO: logger.info("API Key: [CONFIGURED]")
    logger.info("Watsonx API Key: [CONFIGURED]")
    logger.info("Watsonx Project ID: [CONFIGURED]")
    
    return app


# Create the application instance
app = create_application()


@app.get("/", tags=["Health"])
async def root() -> Dict[str, str]:
    """Root endpoint - API information"""
    settings = get_settings()
    return {
        "app": settings.app_name,
        "version": settings.app_version,
        "status": "running",
        "environment": settings.app_env,
        "docs": "/docs" if settings.enable_swagger_docs else "disabled"
    }


@app.get("/health", tags=["Health"])
async def health_check(settings: Settings = Depends(get_settings)) -> Dict[str, Union[str, bool]]:
    """
    Health check endpoint
    Verifies that the application is running and configured correctly
    """
    try:
        # Check if critical configuration is present (Watson NLU is required)
        if not settings.watson_nlu_api_key or not settings.watson_nlu_url:
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail="Missing required Watson NLU configuration"
            )
        
        # Watsonx.ai is optional
        watsonx_configured = bool(settings.watsonx_api_key and settings.watsonx_project_id)
        
        return {
            "status": "healthy",
            "app": settings.app_name,
            "version": settings.app_version,
            "environment": settings.app_env,
            "watsonx_configured": watsonx_configured
        }
    except Exception as e:
        logger.error(f"Health check failed: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Service unhealthy"
        )


@app.get("/config/info", tags=["Configuration"])
async def config_info(settings: Settings = Depends(get_settings)) -> Dict[str, Any]:
    """
    Get non-sensitive configuration information
    NEVER expose API keys or secrets through this endpoint
    """
    return {
        "app_name": settings.app_name,
        "app_version": settings.app_version,
        "environment": settings.app_env,
        "features": {
            "pi_download": settings.enable_pi_download,
            "file_logging": settings.enable_file_logging,
            "swagger_docs": settings.enable_swagger_docs
        },
        "llm_config": {
            "model": settings.watsonx_model,
            "temperature": settings.llm_temperature,
            "max_tokens": settings.llm_max_tokens,
            "top_p": settings.llm_top_p
        },
        "security": {
            "api_key_required": settings.api_key is not None,
            "cors_enabled": True
        }
    }


def extract_text_from_file(file_content: bytes, filename: str) -> str:
    """
    Extract text from uploaded file based on file type
    Supports: TXT, PDF, DOCX, PPTX
    """
    file_extension = filename.lower().split('.')[-1]
    
    try:
        if file_extension == 'txt':
            # Plain text file
            return file_content.decode('utf-8')
        
        elif file_extension == 'pdf':
            # PDF file
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        elif file_extension in ['docx', 'doc']:
            # Word document
            doc_file = io.BytesIO(file_content)
            doc = docx.Document(doc_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text

        elif file_extension == 'pptx':
            # PowerPoint document
            ppt_file = io.BytesIO(file_content)
            presentation = Presentation(ppt_file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text

        elif file_extension == 'ppt':
            # Legacy binary PPT format is not reliably supported by python-pptx
            raise ValueError(
                "Legacy .ppt format is not supported for text extraction. "
                "Please re-save the file as .pptx and upload again."
            )
        
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from file: {str(e)}"
        )


def analyze_project_and_generate_pricing(text: str) -> Dict[str, Any]:
    """
    Analyze project report text and generate pricing estimate
    
    NOTE: This is a simplified demo version. In production, this would:
    - Use Watson NLU for entity extraction and sentiment analysis
    - Use Watsonx.ai LLM for intelligent project analysis
    - Access real pricing databases and historical data
    - Consider location, regulations, and market conditions
    """
    text_lower = text.lower()
    
    # Check for missing critical information
    missing_criteria = []
    warnings = []
    
    # Check for capacity information
    capacity_kw = 100  # Default capacity
    capacity_matches = re.findall(r'(\d+)\s*(?:kw|kilowatt)', text_lower)
    if capacity_matches:
        capacity_kw = int(capacity_matches[0])
    else:
        missing_criteria.append("System capacity (kW)")
        warnings.append("⚠️ Capacity not specified - using default 100kW estimate")
    
    # Check for head height (more flexible pattern)
    head_found = bool(re.search(r'head.*\d+|height.*\d+.*(?:m|meter|metre)|drop.*\d+', text_lower))
    if not head_found:
        missing_criteria.append("Head height (vertical drop)")
        warnings.append("⚠️ Head height not specified - may affect turbine selection")
    
    # Check for flow rate (more flexible pattern)
    flow_found = bool(re.search(r'flow.*\d+|discharge.*\d+|\d+.*(?:m3|cubic|liter|l/s|gpm)', text_lower))
    if not flow_found:
        missing_criteria.append("Water flow rate")
        warnings.append("⚠️ Flow rate not specified - critical for sizing")
    
    # Check for location
    location_found = any(keyword in text_lower for keyword in ['location', 'site', 'address', 'region'])
    if not location_found:
        missing_criteria.append("Project location")
        warnings.append("⚠️ Location not specified - affects pricing and regulations")
    
    # Check for turbine type
    turbine_found = any(keyword in text_lower for keyword in ['francis', 'pelton', 'kaplan', 'crossflow', 'turbine type'])
    if not turbine_found:
        missing_criteria.append("Turbine type preference")
    
    # Detect project type based on keywords
    project_type = "Standard Hydro Turbine"
    if "micro" in text_lower or "small" in text_lower:
        project_type = "Micro Hydro System"
    elif "large" in text_lower or "industrial" in text_lower:
        project_type = "Industrial Hydro Plant"
    
    # Calculate pricing based on capacity and project type
    base_price_per_kw = 1500  # Standard rate
    if "micro" in text_lower:
        base_price_per_kw = 2000  # Higher per-kW cost for smaller systems
    elif "industrial" in text_lower:
        base_price_per_kw = 1200  # Lower per-kW cost for larger systems
    
    # Calculate cost components
    equipment_cost = capacity_kw * base_price_per_kw
    installation_cost = equipment_cost * 0.3  # 30% of equipment cost
    engineering_cost = equipment_cost * 0.15  # 15% of equipment cost
    total_cost = equipment_cost + installation_cost + engineering_cost
    
    # Add realistic variation (±5%)
    variation = random.uniform(0.95, 1.05)
    total_cost = int(total_cost * variation)
    
    # Determine confidence level based on missing criteria
    if len(missing_criteria) == 0:
        confidence_level = "High"
    elif len(missing_criteria) <= 2:
        confidence_level = "Medium"
    else:
        confidence_level = "Low"
    
    # Build notes list
    notes = [
        "This is a preliminary estimate based on document analysis",
        "Final pricing requires detailed site assessment",
        "Prices may vary based on location and specific requirements"
    ]
    
    # Add warnings to notes
    notes.extend(warnings)
    
    # Return structured pricing analysis
    result = {
        "project_type": project_type,
        "estimated_capacity_kw": capacity_kw,
        "pricing_breakdown": {
            "equipment_cost": f"${equipment_cost:,.2f}",
            "installation_cost": f"${installation_cost:,.2f}",
            "engineering_cost": f"${engineering_cost:,.2f}",
            "total_estimated_cost": f"${total_cost:,.2f}"
        },
        "timeline_estimate": "12-18 months",
        "confidence_level": confidence_level,
        "notes": notes
    }
    
    # Add missing criteria if any
    if missing_criteria:
        result["missing_criteria"] = missing_criteria
        result["data_completeness"] = f"{max(0, 100 - len(missing_criteria) * 20)}%"
    else:
        result["data_completeness"] = "100%"
    
    return result


@app.post("/api/analyze-project", tags=["Analysis"])
async def analyze_project_report(
    file: UploadFile = File(...),
    settings: Settings = Depends(get_settings)
) -> Dict[str, Any]:
    """
    Upload a project report and receive an AI-powered pricing estimate
    
    Accepts: TXT, PDF, DOCX, PPTX, PPT files (max 10MB)
    Returns: Detailed pricing breakdown and project analysis
    """
    # Validate filename
    if not file.filename:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Filename is required"
        )
    
    # Validate file type
    allowed_extensions = ['txt', 'pdf', 'docx', 'doc', 'pptx', 'ppt']
    file_extension = file.filename.lower().split('.')[-1]
    
    if file_extension not in allowed_extensions:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}"
        )
    
    # Read and validate file size (max 10MB)
    file_content = await file.read()
    max_size = 10 * 1024 * 1024  # 10MB
    if len(file_content) > max_size:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="File size exceeds 10MB limit"
        )
    
    logger.info(f"Processing file: {file.filename} ({len(file_content)} bytes)")
    
    try:
        # Extract text from uploaded file
        text = extract_text_from_file(file_content, file.filename)
        
        # Validate extracted text
        if not text or len(text.strip()) < 50:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="File appears to be empty or contains insufficient text"
            )
        
        # Analyze project and generate pricing estimate
        analysis_result = analyze_project_and_generate_pricing(text)
        
        # Return comprehensive analysis results
        return {
            "success": True,
            "filename": file.filename,
            "file_size_bytes": len(file_content),
            "text_length": len(text),
            "analysis": analysis_result,
            "timestamp": datetime.utcnow().isoformat() + "Z"
        }
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing project report: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to analyze project report: {str(e)}"
        )


@app.post("/api/generate-proforma-invoice", tags=["Analysis"])
async def generate_proforma_invoice(payload: Dict[str, Any]) -> StreamingResponse:
    """
    Generate a Proforma Invoice Word document from analysis result.
    The document structure follows the provided quotation sample style.
    """
    try:
        analysis = payload.get("analysis") or {}
        pricing = analysis.get("pricing_breakdown") or {}
        turbine_meta = payload.get("turbine_meta") or {}
        source_filename = payload.get("filename") or "inquiry.txt"

        project_type = analysis.get("project_type", "Hydro Turbine Project")
        capacity_kw = analysis.get("estimated_capacity_kw", "N/A")
        timeline = analysis.get("timeline_estimate", "TBD")
        confidence = analysis.get("confidence_level", "N/A")
        turbine_name = turbine_meta.get("name", "Hydro Turbine")
        qty = "1"

        def to_amount(value: Any) -> float:
            s = str(value or "0").replace("$", "").replace(",", "").strip()
            try:
                return float(s)
            except ValueError:
                return 0.0

        # Main-equipment costing base: use equipment_cost as pseudo database base.
        equipment_base = to_amount(pricing.get("equipment_cost", "$0"))
        if equipment_base <= 0:
            # Fallback if upstream data is absent: 70% of total as equipment base.
            equipment_base = to_amount(pricing.get("total_estimated_cost", "$0")) * 0.7

        # Split main equipment into 5 mandatory items.
        # This keeps the quotation structure stable and matches HydroQuote AI scope definition.
        turbine_cost = round(equipment_base * 0.46, 2)
        generator_avr_cost = round(equipment_base * 0.29, 2)
        governor_cost = round(equipment_base * 0.11, 2)
        inlet_valve_cost = round(equipment_base * 0.08, 2)
        automation_cost = round(max(0.0, equipment_base - turbine_cost - generator_avr_cost - governor_cost - inlet_valve_cost), 2)
        main_total = round(turbine_cost + generator_avr_cost + governor_cost + inlet_valve_cost + automation_cost, 2)

        installation_cost = to_amount(pricing.get("installation_cost", "$0"))
        engineering_cost = to_amount(pricing.get("engineering_cost", "$0"))
        grand_total = to_amount(pricing.get("total_estimated_cost", "$0"))

        doc = docx.Document()

        doc.add_paragraph(f"Date: {datetime.utcnow().strftime('%Y-%m-%d')}")
        title = doc.add_paragraph("Proforma invoice")
        if title.runs:
            title.runs[0].bold = True
            title.runs[0].font.size = Pt(16)

        doc.add_paragraph("")
        doc.add_paragraph("Requisites of the Supplier:")
        doc.add_paragraph("Name of supplier: HydroQuote AI Co., Ltd.")
        doc.add_paragraph("Registration number: N/A")
        doc.add_paragraph("Website: https://hydroquote-ai.example.com")
        doc.add_paragraph("Address: Industrial Equipment Zone")
        doc.add_paragraph("E-mail: support@hydroquote-ai.com")
        doc.add_paragraph("Tel: +1 (800) 888-0000")
        doc.add_paragraph("Bank details: Available upon confirmation")
        doc.add_paragraph("")
        doc.add_paragraph("EQUIPMENT:")

        table = doc.add_table(rows=1, cols=5)
        table.style = "Table Grid"
        headers = ["Description", "HS CODE", "Qty", "Unit price", "Total(USD)"]
        for i, h in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = h
            if cell.paragraphs and cell.paragraphs[0].runs:
                cell.paragraphs[0].runs[0].bold = True

        rows = [
            [f"Hydro Turbine ({turbine_name})", "8410120000", qty, f"{turbine_cost:,.2f}", f"{turbine_cost:,.2f}"],
            ["Generator and Exciter / AVR", "8502390010", qty, f"{generator_avr_cost:,.2f}", f"{generator_avr_cost:,.2f}"],
            ["Speed Governor", "9032899090", qty, f"{governor_cost:,.2f}", f"{governor_cost:,.2f}"],
            ["Inlet Valve", "8481804090", qty, f"{inlet_valve_cost:,.2f}", f"{inlet_valve_cost:,.2f}"],
            ["Automation System (Standard PLC Automation)", "8537109090", qty, f"{automation_cost:,.2f}", f"{automation_cost:,.2f}"],
            ["MAIN EQUIPMENT TOTAL", "", "", "", f"{main_total:,.2f}"],
            ["Installation service (excluded from main equipment scope)", "9954990000", qty, f"{installation_cost:,.2f}", f"{installation_cost:,.2f}"],
            ["Engineering service (excluded from main equipment scope)", "9983990000", qty, f"{engineering_cost:,.2f}", f"{engineering_cost:,.2f}"],
            ["GRAND TOTAL", "", "", "", f"{grand_total:,.2f}"]
        ]
        for row_data in rows:
            cells = table.add_row().cells
            for i, value in enumerate(row_data):
                cells[i].text = str(value)

        doc.add_paragraph("")
        doc.add_paragraph(f"Project source file: {source_filename}")
        doc.add_paragraph(f"Selected turbine: {turbine_name}")
        doc.add_paragraph(f"Delivery time: {timeline}")
        doc.add_paragraph("Delivery term (incoterm): FOB / CIF (to be confirmed)")
        doc.add_paragraph("Guarantee: 12 months after commissioning")
        doc.add_paragraph("Country of manufacturing: China")
        doc.add_paragraph(f"Year of manufacturing: {datetime.utcnow().year}")
        doc.add_paragraph(f"Analysis confidence: {confidence}")
        doc.add_paragraph("")
        doc.add_paragraph("Technical Selection Summary")
        doc.add_paragraph(f"Project Summary: {project_type}, estimated capacity {capacity_kw} kW")
        doc.add_paragraph(f"Recommended Turbine Type: {turbine_name}")
        doc.add_paragraph("Unit Configuration: Single-unit baseline (adjustable per customer confirmation)")
        doc.add_paragraph("Main Equipment Scope: Hydro Turbine; Generator and Exciter / AVR; Speed Governor; Inlet Valve; Automation System")
        doc.add_paragraph("Automation Scope: Standard PLC Automation (local HMI, auto start/stop, governor interface, AVR interface, inlet valve interlock, alarm monitoring, unit status display, basic data logging)")
        doc.add_paragraph("Optional Upgrade (not included in main equipment): Advanced SCADA; Remote Monitoring; Central Control Room; Dispatch Integration")
        doc.add_paragraph("Technical Assumptions: preliminary hydraulic/electrical parameters based on provided input and standard selection rules")
        doc.add_paragraph("Risk Notes: transformer, switchgear, penstock, civil works, international freight, and installation service are not included in main equipment scope")
        readiness = "Firm Ready" if confidence == "High" else "Preliminary Ready"
        doc.add_paragraph(f"Quotation Readiness: {readiness}")

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        out_name = f"Proforma_Invoice_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.docx"

        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f'attachment; filename="{out_name}"'}
        )
    except Exception as e:
        logger.error(f"Error generating proforma invoice: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to generate proforma invoice: {str(e)}"
        )


@app.exception_handler(Exception)
async def global_exception_handler(request, exc):
    """
    Global exception handler
    Ensures sensitive information is never leaked in error messages
    """
    logger.error(f"Unhandled exception: {str(exc)}", exc_info=True)
    
    # In production, return generic error message
    settings = get_settings()
    if settings.is_production():
        return JSONResponse(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            content={"detail": "An internal error occurred"}
        )
    
    # In development, provide more details (but still no secrets)
    return JSONResponse(
        status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
        content={
            "detail": "An internal error occurred",
            "error_type": type(exc).__name__,
            "error_message": str(exc)
        }
    )


if __name__ == "__main__":
    import uvicorn
    settings = get_settings()
    
    uvicorn.run(
        "app.main:app",
        host="0.0.0.0",
        port=settings.api_port,
        reload=settings.is_development(),
        log_level=settings.log_level.lower()
    )

# Made with Bob
